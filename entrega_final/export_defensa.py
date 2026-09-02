#!/usr/bin/env python3
"""Exporta el deck HTML a PowerPoint y los PDF de defensa/memoria a Word."""
from __future__ import annotations

import re
import subprocess
import sys
from html import unescape
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "entrega_final" / "deck_defensa_20min.html"
OUT_PPTX = ROOT / "entrega_final" / "TalentPact_Defensa_TFM_20min.pptx"
PNG_DIR = Path("/tmp/talentpact_deck_png")
CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"

CHAPTERS = [
    ("00", "Resumen ejecutivo", "00_resumen_ejecutivo.md"),
    ("M", "Objetivos y metodología", "00_objetivos_metodologia.md"),
    ("01", "Concepto de negocio", "01_concepto.md"),
    ("02", "Estudio de mercado", "02_mercado.md"),
    ("03", "Modelo de negocio", "03_modelo_negocio.md"),
    ("04", "Plan financiero", "04_plan_financiero.md"),
    ("05", "Marketing y ventas", "05_marketing_ventas.md"),
    ("06", "Tecnología y producto", "06_tecnologia_producto.md"),
    ("07", "Regulación y compliance", "07_regulacion_compliance.md"),
    ("08", "Riesgos y contingencias", "08_riesgos.md"),
    ("09", "Conclusiones y limitaciones", "09_conclusiones.md"),
]


def html_to_text(s: str) -> str:
    s = re.sub(r"<br\s*/?>", "\n", s, flags=re.I)
    s = re.sub(r"</p>", "\n", s, flags=re.I)
    s = re.sub(r"<[^>]+>", "", s)
    return unescape(re.sub(r"\n{3,}", "\n\n", s)).strip()


def parse_deck(html: str):
    slides = []
    for m in re.finditer(
        r'<section class="slide[^"]*"([^>]*)>(.*?)</section>',
        html,
        re.S,
    ):
        attrs, body = m.group(1), m.group(2)
        sec = re.search(r'data-sec="([^"]*)"', attrs)
        notes_m = re.search(r'<aside class="notes"[^>]*>(.*?)</aside>', body, re.S)
        title_m = re.search(r"<h1[^>]*>(.*?)</h1>|<h2[^>]*>(.*?)</h2>", body, re.S)
        title = ""
        if title_m:
            title = html_to_text(title_m.group(1) or title_m.group(2) or "")
        notes = html_to_text(notes_m.group(1)) if notes_m else ""
        slides.append(
            {
                "sec": (sec.group(1) if sec else "").strip(),
                "title": title,
                "notes": notes,
            }
        )
    return slides


def screenshot_slides(n: int) -> list[Path]:
    PNG_DIR.mkdir(parents=True, exist_ok=True)
    paths = []
    uri = DECK.resolve().as_uri()
    for i in range(1, n + 1):
        out = PNG_DIR / f"s{i:02d}.png"
        url = f"{uri}?export#{i}"
        cmd = [
            CHROME,
            "--headless=new",
            "--disable-gpu",
            "--hide-scrollbars",
            "--force-device-scale-factor=2",
            "--window-size=1280,720",
            "--virtual-time-budget=4000",
            f"--screenshot={out}",
            url,
        ]
        r = subprocess.run(cmd, capture_output=True, text=True)
        if r.returncode != 0:
            print(r.stderr[-500:], file=sys.stderr)
            raise SystemExit(f"Chrome falló en la diapositiva {i}")
        if not out.exists() or out.stat().st_size < 1000:
            raise SystemExit(f"Captura vacía: {out}")
        paths.append(out)
        print(f"  diapo {i:02d}/{n}  {out.stat().st_size // 1024} KB")
    return paths


def add_hotspot(slide, left_in, top_in, width_in, height_in, url="https://talentpact.es"):
    """Rectángulo transparente clicable encima de la captura."""
    from lxml import etree
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    from pptx.util import Inches

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in),
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    solid = shape._element.spPr.find(qn("a:solidFill"))
    srgb = solid.find(qn("a:srgbClr"))
    alpha = etree.SubElement(srgb, qn("a:alpha"))
    alpha.set("val", "0")
    shape.click_action.hyperlink.address = url


def build_pptx(meta: list[dict], pngs: list[Path]) -> None:
    from pptx import Presentation
    from pptx.util import Emu, Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    # Zonas donde se lee "talentpact.es" (pulgadas, 16:9).
    hotspots = {
        1:  [(0.72, 3.50, 3.50, 0.55)],
        16: [(0.90, 3.85, 7.20, 0.42, "https://talentpact.es/verify.html")],
        17: [(0.72, 0.50, 6.20, 0.42)],
        19: [(0.72, 1.72, 5.40, 0.90)],
        24: [(0.72, 5.00, 4.40, 1.05)],
    }
    for idx, (info, png) in enumerate(zip(meta, pngs), start=1):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(png), Emu(0), Emu(0), width=prs.slide_width, height=prs.slide_height
        )
        for spec in hotspots.get(idx, ()):
            url = spec[4] if len(spec) > 4 else "https://talentpact.es"
            add_hotspot(slide, spec[0], spec[1], spec[2], spec[3], url)
        notes = info["notes"] or info["title"]
        tf = slide.notes_slide.notes_text_frame
        tf.text = notes
    prs.save(OUT_PPTX)
    print(f"PPTX → {OUT_PPTX}  ({len(pngs)} diapositivas)")


def md_to_html_body(md: str) -> str:
    import markdown

    return markdown.markdown(
        md,
        extensions=["tables", "fenced_code", "sane_lists"],
        output_format="html5",
    )


def write_docx_from_html(html: str, title: str, dest: Path) -> None:
    from html2docx import html2docx

    wrapped = (
        "<!DOCTYPE html><html><head><meta charset='utf-8'></head><body>"
        + html
        + "</body></html>"
    )
    buf = html2docx(wrapped, title=title)
    dest.write_bytes(buf.getvalue())
    print(f"DOCX → {dest}  ({dest.stat().st_size // 1024} KB)")


def export_memoria() -> None:
    bp = ROOT / "tfm" / "business_plan"
    parts = [
        "<h1>TalentPact — Plan de negocio</h1>",
        "<p>Trabajo Fin de Máster · Fintech, Mercados Financieros y Blockchain · "
        "Universitat de Barcelona · Xavier Griñó e Ivan Sánchez · septiembre 2026</p>",
    ]
    for num, title, fname in CHAPTERS:
        raw = (bp / fname).read_text(encoding="utf-8")
        lines = raw.splitlines()
        if lines and lines[0].startswith("# "):
            raw = "\n".join(lines[1:])
        raw = re.sub(r"\n---\n\s*\*[^*]+\*\s*$", "", raw.strip())
        parts.append(f"<h1>{num}. {title}</h1>")
        parts.append(md_to_html_body(raw))
    dest = ROOT / "tfm" / "TalentPact_TFM_Business_Plan.docx"
    write_docx_from_html("\n".join(parts), "TalentPact — Plan de negocio", dest)


def export_md_pdfs() -> None:
    pairs = [
        (
            ROOT / "entrega_final" / "GUION_DEFENSA_20MIN.md",
            ROOT / "entrega_final" / "GUION_DEFENSA_20MIN.docx",
            "Guion de defensa — 20 minutos",
        ),
        (
            ROOT / "entrega_final" / "QA_DEFENSA.md",
            ROOT / "entrega_final" / "QA_DEFENSA.docx",
            "Q&A de defensa",
        ),
    ]
    for src, dest, title in pairs:
        if not src.exists():
            continue
        write_docx_from_html(md_to_html_body(src.read_text(encoding="utf-8")), title, dest)


def main() -> None:
    html = DECK.read_text(encoding="utf-8")
    meta = parse_deck(html)
    print(f"Deck: {len(meta)} diapositivas")
    print("Capturando diapositivas…")
    pngs = screenshot_slides(len(meta))
    build_pptx(meta, pngs)
    print("Convirtiendo PDFs a Word…")
    export_memoria()
    export_md_pdfs()


if __name__ == "__main__":
    main()
